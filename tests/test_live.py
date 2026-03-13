"""Live integration tests — validate tools against real inputs.

These tests create real files and (optionally) make real network requests.
Mark network tests with @mark("network") so they can be skipped offline.
"""

from pathlib import Path
from stratagem.testing import mark


# ── scrape_url ──────────────────────────────────────────────

class TestScrapeUrlLive:
    @mark("network")
    async def test_scrape_real_url(self):
        from stratagem.tools.web import scrape_url

        result = await scrape_url.handler({"url": "https://example.com"})
        assert "isError" not in result or not result["isError"]
        text = result["content"][0]["text"]
        assert len(text) > 50
        assert "example" in text.lower() or "Example" in text


# ── read_spreadsheet ───────────────────────────────────────

class TestReadSpreadsheetLive:
    async def test_roundtrip_xlsx(self, tmp_dir):
        """Create a real .xlsx then parse it back."""
        from openpyxl import Workbook
        from stratagem.tools.spreadsheet import read_spreadsheet

        # Create test spreadsheet
        wb = Workbook()
        ws = wb.active
        ws.title = "Sales"
        ws.append(["Product", "Q1", "Q2", "Q3"])
        ws.append(["Widget A", 100, 200, 150])
        ws.append(["Widget B", 300, 250, 400])
        ws.append(["Widget C", 50, 75, 60])
        xlsx_path = str(tmp_dir / "test_data.xlsx")
        wb.save(xlsx_path)

        # Parse it back
        result = await read_spreadsheet.handler({"file_path": xlsx_path})
        assert "isError" not in result or not result["isError"]
        text = result["content"][0]["text"]
        assert "Product" in text
        assert "Widget A" in text
        assert "Sales" in text


# ── create_report (markdown) ──────────────────────────────

class TestCreateReportLive:
    async def test_markdown_report(self, tmp_dir):
        from stratagem.tools.reports import create_report

        out_path = str(tmp_dir / "test_report.md")
        result = await create_report.handler({
            "title": "Test Research Report",
            "sections": [
                {"heading": "Introduction", "content": "This is the intro section.", "level": 2},
                {"heading": "Findings", "content": "- Finding 1\n- Finding 2\n- Finding 3", "level": 2},
                {"heading": "Conclusion", "content": "Summary of conclusions.", "level": 2},
            ],
            "format": "markdown",
            "output_path": out_path,
            "metadata": {"author": "Test Suite", "date": "2026-03-11"},
        })

        assert "isError" not in result or not result["isError"]
        assert Path(out_path).exists()
        content = Path(out_path).read_text()
        assert "# Test Research Report" in content
        assert "Introduction" in content
        assert "Finding 1" in content
        assert "Test Suite" in content

    async def test_docx_report(self, tmp_dir):
        """Test new DOCX output format."""
        from stratagem.tools.reports import create_report

        out_path = str(tmp_dir / "test_report.docx")
        result = await create_report.handler({
            "title": "DOCX Test Report",
            "sections": [
                {"heading": "Overview", "content": "This tests DOCX generation.", "level": 2},
                {"heading": "Bullet Points", "content": "- Item one\n- Item two\n- Item three", "level": 2},
                {"heading": "Details", "content": "Paragraph with details.\nAnother line.", "level": 3},
            ],
            "format": "docx",
            "output_path": out_path,
            "metadata": {"author": "Stratagem Test", "date": "2026-03-11"},
        })

        assert "isError" not in result or not result["isError"]
        assert Path(out_path).exists()
        assert Path(out_path).stat().st_size > 0

        # Verify it's a valid DOCX by reading it back
        from docx import Document
        doc = Document(out_path)
        full_text = "\n".join(p.text for p in doc.paragraphs)
        assert "DOCX Test Report" in full_text
        assert "Item one" in full_text

    async def test_html_report(self, tmp_dir):
        from stratagem.tools.reports import create_report

        out_path = str(tmp_dir / "test_report.html")
        result = await create_report.handler({
            "title": "HTML Test",
            "sections": [{"heading": "Sec1", "content": "Content here.", "level": 2}],
            "format": "html",
            "output_path": out_path,
        })

        assert "isError" not in result or not result["isError"]
        content = Path(out_path).read_text()
        assert "<h1>" in content
        assert "HTML Test" in content


# ── PPTX roundtrip ────────────────────────────────────────

class TestPptxRoundtrip:
    async def test_create_and_read_pptx(self, tmp_dir):
        from stratagem.tools.reports import create_report
        from stratagem.tools.presentation import read_pptx

        pptx_path = str(tmp_dir / "test_slides.pptx")

        # Create PPTX
        create_result = await create_report.handler({
            "title": "Test Presentation",
            "sections": [
                {"heading": "Slide One", "content": "First slide content.", "level": 2},
                {"heading": "Slide Two", "content": "Second slide content.", "level": 2},
            ],
            "format": "pptx",
            "output_path": pptx_path,
            "metadata": {"author": "Test"},
        })

        assert "isError" not in create_result or not create_result["isError"]
        assert Path(pptx_path).exists()

        # Read it back
        read_result = await read_pptx.handler({"file_path": pptx_path})
        assert "isError" not in read_result or not read_result["isError"]
        text = read_result["content"][0]["text"]
        assert "Test Presentation" in text
        assert "Slide One" in text


# ── extract_images ────────────────────────────────────────

class TestExtractImagesLive:
    async def test_extract_from_pptx(self, tmp_dir):
        """Create a PPTX with an image, then extract it."""
        from pptx import Presentation
        from pptx.util import Inches
        from PIL import Image
        from stratagem.tools.images import extract_images

        # Create a small test image
        img_path = tmp_dir / "test_img.png"
        img = Image.new("RGB", (100, 100), color="blue")
        img.save(str(img_path))

        # Create PPTX with the image
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        slide.shapes.add_picture(str(img_path), Inches(1), Inches(1), Inches(2), Inches(2))
        pptx_path = str(tmp_dir / "with_image.pptx")
        prs.save(pptx_path)

        # Extract images
        extract_dir = str(tmp_dir / "extracted")
        result = await extract_images.handler({"file_path": pptx_path, "output_dir": extract_dir})
        assert "isError" not in result or not result["isError"]
        text = result["content"][0]["text"]
        assert "image" in text.lower() or "extracted" in text.lower() or "1" in text


# ── SEC EDGAR ─────────────────────────────────────────────

class TestSecEdgarLive:
    @mark("network")
    async def test_search_aapl_filings(self):
        from stratagem.tools.sec_edgar import search_sec_filings

        result = await search_sec_filings.handler({
            "ticker": "AAPL",
            "form_type": "10-K",
            "limit": 3,
        })
        assert "isError" not in result or not result["isError"]
        text = result["content"][0]["text"]
        assert "AAPL" in text
        assert "10-K" in text

    async def test_empty_filings_no_crash(self):
        """Verify querying a bogus ticker returns error, not crash."""
        from stratagem.tools.sec_edgar import download_sec_filing

        result = await download_sec_filing.handler({
            "ticker": "ZZZZZZ999",
            "form_type": "10-K",
            "filing_index": 0,
        })
        # Should return an error, not crash
        assert result.get("isError") is True or "Error" in result["content"][0]["text"] or "not found" in result["content"][0]["text"].lower() or "Failed" in result["content"][0]["text"]
