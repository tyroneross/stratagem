"""PowerPoint tools — read and create .pptx files."""

from typing import Any
from pathlib import Path

from claude_agent_sdk import tool


@tool(
    "read_pptx",
    "Read a PowerPoint file and extract text, speaker notes, and optionally images from each slide.",
    {
        "type": "object",
        "properties": {
            "file_path": {"type": "string", "description": "Absolute path to the .pptx file"},
            "extract_images": {"type": "boolean", "description": "Extract embedded images", "default": False},
            "include_notes": {"type": "boolean", "description": "Include speaker notes", "default": True},
        },
        "required": ["file_path"],
    },
)
async def read_pptx(args: dict[str, Any]) -> dict[str, Any]:
    file_path = args["file_path"]
    extract_images = args.get("extract_images", False)
    include_notes = args.get("include_notes", True)

    try:
        from pptx import Presentation
    except ImportError:
        return _error("python-pptx not installed. Run: pip install python-pptx")

    path = Path(file_path)
    if not path.exists():
        return _error(f"File not found: {file_path}")

    try:
        prs = Presentation(file_path)
    except Exception as e:
        return _error(f"Failed to open PPTX: {e}")

    sections: list[str] = [
        f"# Presentation: {path.name}",
        f"**Slides**: {len(prs.slides)}",
        "",
    ]

    for slide_num, slide in enumerate(prs.slides, 1):
        sections.append(f"## Slide {slide_num}")

        # Extract text from shapes
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)

            # Table shapes
            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(cells)
                if rows:
                    texts.append(_rows_to_markdown(rows))

        if texts:
            sections.append("\n".join(texts))
        else:
            sections.append("*(no text content)*")

        # Speaker notes
        if include_notes and slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                sections.append(f"\n**Speaker Notes**: {notes_text}")

        # Images
        if extract_images:
            import base64
            img_count = 0
            for shape in slide.shapes:
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    try:
                        image = shape.image
                        b64 = base64.b64encode(image.blob).decode("ascii")
                        img_count += 1
                        sections.append(f"\n*Image {img_count}: {image.content_type}, {len(image.blob)} bytes*")
                    except Exception:
                        continue
            if img_count == 0:
                sections.append("\n*No images on this slide*")

        sections.append("")

    return {"content": [{"type": "text", "text": "\n".join(sections)}]}


@tool(
    "create_pptx",
    "Create a PowerPoint presentation from structured slide data. Supports title, bullet, table, and image slides.",
    {
        "type": "object",
        "properties": {
            "slides": {
                "type": "array",
                "description": "Array of slide objects with title, content, and layout type",
                "items": {
                    "type": "object",
                    "properties": {
                        "title": {"type": "string"},
                        "content": {"type": "string", "description": "Bullet text (newline-separated) or table data"},
                        "layout": {"type": "string", "enum": ["title", "bullets", "table", "blank"], "default": "bullets"},
                        "notes": {"type": "string", "description": "Speaker notes for this slide"},
                    },
                    "required": ["title"],
                },
            },
            "output_path": {"type": "string", "description": "Where to save the .pptx file"},
            "template": {"type": "string", "description": "Optional path to a template .pptx file"},
        },
        "required": ["slides", "output_path"],
    },
)
async def create_pptx(args: dict[str, Any]) -> dict[str, Any]:
    slides_data = args["slides"]
    output_path = args["output_path"]
    template = args.get("template")

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        return _error("python-pptx not installed. Run: pip install python-pptx")

    # Ensure output directory exists
    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)

    try:
        prs = Presentation(template) if template else Presentation()
    except Exception as e:
        return _error(f"Failed to create presentation: {e}")

    for slide_data in slides_data:
        title = slide_data.get("title", "")
        content = slide_data.get("content", "")
        layout_type = slide_data.get("layout", "bullets")
        notes = slide_data.get("notes", "")

        if layout_type == "title":
            layout = prs.slide_layouts[0]  # Title Slide
            slide = prs.slides.add_slide(layout)
            if slide.placeholders[0]:
                slide.placeholders[0].text = title
            if len(slide.placeholders) > 1 and content:
                slide.placeholders[1].text = content

        elif layout_type == "bullets":
            layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(layout)
            if slide.placeholders[0]:
                slide.placeholders[0].text = title
            if len(slide.placeholders) > 1 and content:
                tf = slide.placeholders[1].text_frame
                tf.text = ""
                for i, line in enumerate(content.split("\n")):
                    line = line.strip().lstrip("- •")
                    if not line:
                        continue
                    if i == 0:
                        tf.text = line
                    else:
                        p = tf.add_paragraph()
                        p.text = line

        elif layout_type == "table":
            layout = prs.slide_layouts[5]  # Blank
            slide = prs.slides.add_slide(layout)
            # Add title as text box
            from pptx.util import Inches, Pt
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
            txBox.text_frame.text = title
            txBox.text_frame.paragraphs[0].font.size = Pt(24)
            txBox.text_frame.paragraphs[0].font.bold = True

            # Parse table from content (pipe-delimited or CSV)
            rows = _parse_table_content(content)
            if rows:
                row_count = len(rows)
                col_count = max(len(r) for r in rows)
                tbl = slide.shapes.add_table(
                    row_count, col_count, Inches(0.5), Inches(1.5), Inches(9), Inches(4.5)
                ).table
                for r_idx, row in enumerate(rows):
                    for c_idx, cell_val in enumerate(row):
                        if c_idx < col_count:
                            tbl.cell(r_idx, c_idx).text = cell_val

        else:  # blank
            layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[5]
            slide = prs.slides.add_slide(layout)

        # Add speaker notes
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    try:
        prs.save(output_path)
    except Exception as e:
        return _error(f"Failed to save presentation: {e}")

    return {"content": [{"type": "text", "text": f"Created presentation: {output_path} ({len(slides_data)} slides)"}]}


def _parse_table_content(content: str) -> list[list[str]]:
    """Parse table content from pipe-delimited or CSV format."""
    rows = []
    for line in content.strip().split("\n"):
        line = line.strip()
        if not line or line.startswith("|---") or set(line) <= {"|", "-", " "}:
            continue
        if "|" in line:
            cells = [c.strip() for c in line.strip("|").split("|")]
            rows.append(cells)
        elif "," in line:
            import csv as csv_mod
            import io
            reader = csv_mod.reader(io.StringIO(line))
            for row in reader:
                rows.append(row)
                break
    return rows


def _rows_to_markdown(rows: list[list[str]]) -> str:
    if not rows:
        return ""
    col_count = max(len(r) for r in rows)
    lines = []
    for i, row in enumerate(rows):
        padded = row + [""] * (col_count - len(row))
        lines.append("| " + " | ".join(padded) + " |")
        if i == 0:
            lines.append("| " + " | ".join(["---"] * col_count) + " |")
    return "\n".join(lines)


def _error(msg: str) -> dict[str, Any]:
    return {"content": [{"type": "text", "text": f"Error: {msg}"}], "isError": True}
