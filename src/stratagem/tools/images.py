"""Image extraction tool — extract images from PDFs, PPTX, and DOCX files."""

from typing import Any
from pathlib import Path
import base64

from claude_agent_sdk import tool


@tool(
    "extract_images",
    "Extract embedded images from PDF, PPTX, or DOCX files. Returns image metadata and base64 data.",
    {
        "type": "object",
        "properties": {
            "file_path": {"type": "string", "description": "Absolute path to the document"},
            "output_dir": {"type": "string", "description": "Optional directory to save images to"},
        },
        "required": ["file_path"],
    },
)
async def extract_images(args: dict[str, Any]) -> dict[str, Any]:
    file_path = args["file_path"]
    output_dir = args.get("output_dir")

    path = Path(file_path)
    if not path.exists():
        return _error(f"File not found: {file_path}")

    ext = path.suffix.lower()

    if ext == ".pdf":
        images = _extract_from_pdf(path)
    elif ext == ".pptx":
        images = _extract_from_pptx(path)
    elif ext == ".docx":
        images = _extract_from_docx(path)
    else:
        return _error(f"Unsupported format: {ext}. Supported: .pdf, .pptx, .docx")

    if isinstance(images, dict) and images.get("isError"):
        return images

    # Save to disk if output_dir specified
    if output_dir and images:
        out_path = Path(output_dir)
        out_path.mkdir(parents=True, exist_ok=True)
        for img in images:
            if img.get("data"):
                img_file = out_path / img["name"]
                img_file.write_bytes(base64.b64decode(img["data"]))
                img["saved_to"] = str(img_file)

    # Build response
    sections = [
        f"# Images from: {path.name}",
        f"**Total images**: {len(images)}",
        "",
    ]

    for img in images:
        sections.append(f"- **{img['name']}**: {img.get('content_type', 'unknown')} ({img.get('size', 0)} bytes)")
        if img.get("source"):
            sections.append(f"  Source: {img['source']}")
        if img.get("saved_to"):
            sections.append(f"  Saved: {img['saved_to']}")

    if not images:
        sections.append("No images found.")

    return {"content": [{"type": "text", "text": "\n".join(sections)}]}


def _extract_from_pdf(path: Path) -> list[dict]:
    try:
        from pypdf import PdfReader
    except ImportError:
        return _error("pypdf not installed")

    images = []
    reader = PdfReader(str(path))

    for page_num, page in enumerate(reader.pages, 1):
        if hasattr(page, "images"):
            for img_idx, image in enumerate(page.images):
                try:
                    name = getattr(image, "name", f"page{page_num}_img{img_idx}.png")
                    images.append({
                        "name": name,
                        "content_type": _guess_content_type(name),
                        "size": len(image.data),
                        "source": f"Page {page_num}",
                        "data": base64.b64encode(image.data).decode("ascii"),
                    })
                except Exception:
                    continue

    return images


def _extract_from_pptx(path: Path) -> list[dict]:
    try:
        from pptx import Presentation
    except ImportError:
        return _error("python-pptx not installed")

    images = []
    prs = Presentation(str(path))

    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture
                try:
                    image = shape.image
                    ext = image.content_type.split("/")[-1] if image.content_type else "png"
                    name = f"slide{slide_num}_{shape.name}.{ext}"
                    images.append({
                        "name": name,
                        "content_type": image.content_type,
                        "size": len(image.blob),
                        "source": f"Slide {slide_num}",
                        "data": base64.b64encode(image.blob).decode("ascii"),
                    })
                except Exception:
                    continue

    return images


def _extract_from_docx(path: Path) -> list[dict]:
    """Extract images from DOCX using zipfile (DOCX is a ZIP archive)."""
    import zipfile

    images = []

    try:
        with zipfile.ZipFile(str(path), "r") as z:
            for name in z.namelist():
                if name.startswith("word/media/"):
                    data = z.read(name)
                    file_name = Path(name).name
                    images.append({
                        "name": file_name,
                        "content_type": _guess_content_type(file_name),
                        "size": len(data),
                        "source": "word/media/",
                        "data": base64.b64encode(data).decode("ascii"),
                    })
    except zipfile.BadZipFile:
        return _error("Invalid DOCX file (not a valid ZIP archive)")
    except Exception as e:
        return _error(f"DOCX extraction error: {e}")

    return images


def _guess_content_type(filename: str) -> str:
    ext = Path(filename).suffix.lower()
    return {
        ".png": "image/png",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".gif": "image/gif",
        ".bmp": "image/bmp",
        ".tiff": "image/tiff",
        ".svg": "image/svg+xml",
        ".webp": "image/webp",
        ".emf": "image/x-emf",
        ".wmf": "image/x-wmf",
    }.get(ext, "application/octet-stream")


def _error(msg: str) -> dict[str, Any]:
    return {"content": [{"type": "text", "text": f"Error: {msg}"}], "isError": True}
